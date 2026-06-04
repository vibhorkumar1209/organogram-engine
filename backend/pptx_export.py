"""
PPTX export — org-chart style.

Card design: compact light-gray rounded box, coloured left-accent bar, dark text.
  Name (bold) / Designation / ─── / location / LinkedIn
Layout: dept-head centred at top → column heads below → members stacked.
Connector style: vertical stem → horizontal bar → drops (T-bar).
"""
from __future__ import annotations

import io
from datetime import datetime
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt

# ── Slide geometry ─────────────────────────────────────────────────────────────

SLIDE_W  = Inches(13.333)
SLIDE_H  = Inches(7.5)
HEADER_H = Inches(1.05)

# ── Card geometry ─────────────────────────────────────────────────────────────
#
# CARD_H tuned so ~4-5 member cards fit below the column head on one slide:
#   Available = 7.5 - 1.05(header) - 0.20(gap) - 0.82(dept-head)
#               - 0.36(connectors) - 0.82(col-head) - 0.06(gap) - 0.12(bottom)
#             ≈ 4.07"
#   Cards/col = floor(4.07 / (0.82 + 0.06)) = 4
#
CARD_H     = Inches(0.82)   # compact — was 1.32"
HEAD_W     = Inches(2.10)   # dept-head card is slightly wider
COL_GAP    = Inches(0.26)   # horizontal gap between columns
CARD_V_GAP = Inches(0.06)   # vertical gap between stacked cards
CONN_H     = Inches(0.18)   # height reserved for connector lines
SIDE_PAD   = Inches(0.45)
TREE_TOP   = int(HEADER_H) + int(Inches(0.18))
MAX_COLS   = 6              # max columns per slide before paginating

# Accent-bar width (left strip with dept / branch colour)
ACCENT_W = int(Inches(0.055))

# ── Card colours (readable at all sizes) ─────────────────────────────────────

_CARD_BG   = RGBColor(0xed, 0xf1, 0xf6)   # light blue-gray fill
_CARD_BG_H = RGBColor(0xde, 0xe6, 0xf3)   # slightly darker for dept-head card
_CARD_BDR  = RGBColor(0x8e, 0xa6, 0xc0)   # visible gray border
_TXT_NAME  = RGBColor(0x0f, 0x1d, 0x2c)   # near-black bold  ← always inside card
_TXT_TITLE = RGBColor(0x26, 0x40, 0x5c)   # dark slate-blue
_TXT_BODY  = RGBColor(0x38, 0x52, 0x6c)   # dark medium-gray
_TXT_LINK  = RGBColor(0x16, 0x65, 0xbc)   # blue hyperlink
_DIVIDER   = RGBColor(0xb0, 0xc4, 0xd8)

# Cycling accent palette for sub-depts
_PALETTE: list[tuple[int, int, int]] = [
    (0x1e, 0x6e, 0xb8),
    (0x16, 0x7a, 0x56),
    (0x7a, 0x1c, 0x94),
    (0xb0, 0x24, 0x24),
    (0xc8, 0x6c, 0x1a),
    (0x0e, 0x68, 0x82),
]

# ── Colour helpers ─────────────────────────────────────────────────────────────

def _luminance(r: int, g: int, b: int) -> float:
    return 0.299*r + 0.587*g + 0.114*b

def _on_bg(r: int, g: int, b: int) -> RGBColor:
    return RGBColor(0xff,0xff,0xff) if _luminance(r,g,b) < 130 else RGBColor(0x0f,0x1d,0x2c)

def _hex_to_rgb(h: str) -> tuple[int, int, int]:
    h = h.lstrip("#")
    if len(h) != 6:
        return (0x1e, 0x6e, 0xb8)
    return (int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def _rgb(r: int, g: int, b: int) -> RGBColor:
    return RGBColor(r, g, b)


# ── Low-level drawing helpers ─────────────────────────────────────────────────

def _add_rect(slide, left, top, width, height,
              fill_rgb: RGBColor,
              line_rgb: RGBColor | None = None,
              line_width: float = 0,
              corner_radius: int = 0):
    """Always requires fill_rgb — never calls fill.background() which paints white."""
    from pptx.oxml.ns import qn
    from lxml import etree

    shape = slide.shapes.add_shape(1, int(left), int(top), int(width), int(height))

    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb

    shape.line.width = int(Pt(line_width)) if line_width else 0
    if line_rgb:
        shape.line.color.rgb = line_rgb
    else:
        shape.line.fill.background()   # no border line

    if corner_radius:
        sp   = shape._element
        spPr = sp.find(qn("p:spPr"))
        pg   = spPr.find(qn("a:prstGeom"))
        if pg is not None:
            spPr.remove(pg)
        ng  = etree.SubElement(spPr, qn("a:prstGeom"), attrib={"prst": "roundRect"})
        av  = etree.SubElement(ng, qn("a:avLst"))
        adj = min(50000, int(corner_radius * 100000 // min(width, height)))
        etree.SubElement(av, qn("a:gd"), attrib={"name":"adj","fmla":f"val {adj}"})
    return shape


def _add_textbox(slide, left, top, width, height, text: str,
                 font_size: float, bold: bool = False,
                 color: RGBColor | None = None,
                 align=PP_ALIGN.LEFT, wrap: bool = True):
    """
    Fixed-size textbox — auto_size=NONE so text never escapes the declared bounds.
    wrap=True (default) prevents horizontal overflow beyond card edges.
    """
    tb = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE   # ← critical: no shape expansion
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return tb


def _add_hyperlink_box(slide, left, top, width, height,
                       display: str, url: str, font_size: float):
    from pptx.oxml.ns import qn
    from lxml import etree
    tb  = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tf  = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE   # keep within bounds
    p   = tf.paragraphs[0]
    run = p.add_run()
    run.text = display
    run.font.size = Pt(font_size)
    run.font.color.rgb = _TXT_LINK
    try:
        rId = slide.part.relate_to(
            url,
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
            is_external=True)
        rPr = run._r.get_or_add_rPr()
        hl  = etree.SubElement(rPr, qn("a:hlinkClick"))
        hl.set(qn("r:id"), rId)
    except Exception:
        pass


_LP = 9525   # ~0.75 pt in EMU

def _add_line(slide, x1, y1, x2, y2, color: RGBColor):
    if abs(x2 - x1) < _LP:
        _add_rect(slide, x1-_LP//2, min(y1,y2), _LP, max(abs(y2-y1),_LP), fill_rgb=color)
    else:
        _add_rect(slide, min(x1,x2), y1-_LP//2, abs(x2-x1), max(abs(y2-y1),_LP), fill_rgb=color)


# ── Cover slide ───────────────────────────────────────────────────────────────

def _make_cover(prs, company, total_people, dept_count, industry):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=RGBColor(0x08,0x15,0x22))
    _add_rect(slide, 0, 0, Inches(0.22), SLIDE_H, fill_rgb=RGBColor(0x34,0x91,0xE8))
    _add_textbox(slide, Inches(0.55), Inches(2.10), Inches(11.5), Inches(1.6),
                 company, font_size=44, bold=True, color=RGBColor(0xff,0xff,0xff))
    sub = f"Organisational Chart  ·  {total_people:,} people  ·  {dept_count} departments"
    if industry:
        sub += f"  ·  {industry}"
    _add_textbox(slide, Inches(0.55), Inches(3.80), Inches(11.5), Inches(0.55),
                 sub, font_size=13, color=RGBColor(0x7e,0xc8,0xf8))
    _add_textbox(slide, Inches(0.55), Inches(6.65), Inches(6.0), Inches(0.4),
                 f"Generated {datetime.now().strftime('%d %B %Y')}",
                 font_size=9, color=RGBColor(0x44,0x6e,0x88))


# ── Slide header bar ──────────────────────────────────────────────────────────

def _draw_header(slide, dept_label, company, dept_rgb, total_hc, page_sfx=""):
    r, g, b = dept_rgb
    fill = _rgb(r, g, b)
    txt  = _on_bg(r, g, b)
    _add_rect(slide, 0, 0, SLIDE_W, HEADER_H, fill_rgb=fill)
    _add_rect(slide, 0, HEADER_H - int(Inches(0.04)), SLIDE_W, int(Inches(0.04)),
              fill_rgb=RGBColor(0xff,0xff,0xff))
    label = dept_label.upper() + (f"  {page_sfx}" if page_sfx else "")
    _add_textbox(slide, Inches(0.45), Inches(0.14), Inches(9.0), Inches(0.70),
                 label, font_size=22, bold=True, color=txt)
    _add_textbox(slide, Inches(10.0), Inches(0.20), Inches(3.1), Inches(0.55),
                 f"{company}\n{total_hc:,} people", font_size=9, color=txt,
                 align=PP_ALIGN.RIGHT)


# ── Org-chart person card ─────────────────────────────────────────────────────
#
#  ┌▌────────────────────────────────────────┐   ← 0.82" tall
#  │ Alex Vance              (bold 9pt dark) │
#  │ Co-Founder & CEO        (7.5pt slate)   │
#  │ ──────────────────────────────────────  │
#  │ Location: Florida       (6.5pt gray)    │
#  │ linkedin.com/in/alexvance (blue link)   │
#  └─────────────────────────────────────────┘
#
# IMPORTANT:
#  • wrap=True on all text → text wraps at card boundary, never overflows to right
#  • auto_size=NONE        → shape never expands, text clips at bottom if needed
#  • single gray rect fill → no layered transparent shapes
#
def _draw_card(slide, left: int, top: int,
               person: dict,
               accent_rgb: tuple[int, int, int],
               card_w: int | None = None,
               is_head: bool = False) -> None:

    cw = card_w or int(Inches(1.92))
    ch = int(CARD_H)
    bg = _CARD_BG_H if is_head else _CARD_BG

    # ── 1. Single gray card background (one rect, no layers) ─────────────
    _add_rect(slide, left, top, cw, ch,
              fill_rgb=bg,
              line_rgb=_CARD_BDR,
              line_width=0.9,
              corner_radius=4000)

    # ── 2. Left accent bar ────────────────────────────────────────────────
    r, g, b = accent_rgb
    _add_rect(slide, left, top, ACCENT_W, ch,
              fill_rgb=_rgb(r, g, b),
              corner_radius=4000)

    # ── 3. Text content — all inside the card ────────────────────────────
    #  ix: x-start of text  (after accent bar + inner padding)
    #  iw: text width        (card width minus accent bar minus right pad)
    ix  = left + ACCENT_W + int(Inches(0.07))
    iw  = cw - ACCENT_W - int(Inches(0.10))
    # Clamp iw so it's never zero/negative
    iw  = max(int(Inches(0.5)), iw)
    cy  = top + int(Inches(0.06))

    name  = str(person.get("label") or "").strip()
    title = str((person.get("metadata") or {}).get("designation") or "").strip()
    meta  = person.get("metadata") or {}

    # Name — wrap=True so long names wrap rather than overflow right edge
    if name:
        _add_textbox(slide, ix, cy, iw, int(Inches(0.20)),
                     name, font_size=9.0, bold=True, color=_TXT_NAME, wrap=True)
        cy += int(Inches(0.20))

    # Designation — wrap=True
    if title:
        _add_textbox(slide, ix, cy, iw, int(Inches(0.17)),
                     title, font_size=7.5, color=_TXT_TITLE, wrap=True)
        cy += int(Inches(0.17))

    # Thin divider
    cy += int(Inches(0.03))
    _add_rect(slide, ix, cy, iw, max(1, int(Inches(0.013))), fill_rgb=_DIVIDER)
    cy += int(Inches(0.04))

    # Contact detail lines — stop if we'd overflow the card bottom
    line_h  = int(Inches(0.15))
    line_gap = int(Inches(0.155))
    max_y   = top + ch - int(Inches(0.04))

    def _line(text: str, url: str | None = None) -> bool:
        nonlocal cy
        if cy + line_h > max_y:
            return False
        if url:
            _add_hyperlink_box(slide, ix, cy, iw, line_h, text, url, 6.5)
        else:
            _add_textbox(slide, ix, cy, iw, line_h, text, 6.5, color=_TXT_BODY, wrap=True)
        cy += line_gap
        return True

    loc = str(meta.get("city") or meta.get("region") or meta.get("location") or "").strip()[:35]
    li  = str(meta.get("linkedin_url") or meta.get("linkedin") or
              meta.get("LinkedInURL") or "").strip()
    email = str(meta.get("email") or "").strip()

    if loc   and not _line(f"📍 {loc}"): return
    if email and not _line(email): return
    if li:
        disp = li.replace("https://www.","").replace("https://","").replace("http://","")[:40]
        _line(disp, url=li)


# ── Connector lines ───────────────────────────────────────────────────────────

_CONN_CLR = RGBColor(0x8a, 0xa2, 0xba)

def _draw_connectors(slide, parent_cx, parent_bottom,
                     child_cxs: list[int], child_top) -> None:
    if not child_cxs:
        return
    mid_y = parent_bottom + (child_top - parent_bottom) // 2
    if len(child_cxs) == 1:
        _add_line(slide, parent_cx, parent_bottom, parent_cx, child_top, _CONN_CLR)
    else:
        _add_line(slide, parent_cx, parent_bottom, parent_cx, mid_y, _CONN_CLR)
        _add_line(slide, min(child_cxs), mid_y, max(child_cxs), mid_y, _CONN_CLR)
        for cx in child_cxs:
            _add_line(slide, cx, mid_y, cx, child_top, _CONN_CLR)


# ── Department slide(s) ───────────────────────────────────────────────────────

def _make_dept_slides(prs: Presentation, company: str,
                      dept_label: str, dept_color_hex: str,
                      total_hc: int,
                      dept_head: dict | None,
                      columns: list[dict]) -> None:
    dept_rgb = _hex_to_rgb(dept_color_hex)
    blank    = prs.slide_layouts[6]

    avail_w = int(SLIDE_W) - 2 * int(SIDE_PAD)

    def _col_width(n: int) -> int:
        nominal_total = n * int(Inches(1.92)) + (n-1) * int(COL_GAP)
        if nominal_total <= avail_w:
            return int(Inches(1.92))
        return max(int(Inches(1.30)), (avail_w - (n-1)*int(COL_GAP)) // n)

    def _one_slide(slide_cols: list[dict], sfx: str) -> None:
        slide = prs.slides.add_slide(blank)
        _draw_header(slide, dept_label, company, dept_rgb, total_hc, sfx)

        n   = len(slide_cols)
        cw  = _col_width(n)
        gap = int(COL_GAP) if cw == int(Inches(1.92)) else int(Inches(0.15))
        ch  = int(CARD_H)

        total_w   = n * cw + (n-1) * gap
        col_start = (int(SLIDE_W) - total_w) // 2
        col_xs    = [col_start + i*(cw+gap) for i in range(n)]
        col_cxs   = [x + cw//2 for x in col_xs]

        # Dept-head card — centred at top
        head_y   = TREE_TOP
        head_cw  = min(int(HEAD_W), avail_w)
        head_x   = int(SLIDE_W//2) - head_cw//2
        head_btm = head_y + ch

        if dept_head:
            _draw_card(slide, head_x, head_y, dept_head, dept_rgb,
                       card_w=head_cw, is_head=True)

        # Gap for connector lines
        col_head_y = head_btm + int(CONN_H) * 2

        # T-bar connectors: dept head → column heads
        if dept_head and slide_cols:
            _draw_connectors(slide, int(SLIDE_W//2), head_btm, col_cxs, col_head_y)

        # Draw columns
        for i, col in enumerate(slide_cols):
            cx  = col_xs[i]
            acc = col["accent_rgb"]

            if col.get("head"):
                # Column has a named head (sub-dept head or first person in group)
                _draw_card(slide, cx, col_head_y, col["head"], acc, card_w=cw)
                member_y = col_head_y + ch + int(CARD_V_GAP)
            else:
                # Continuation column — no head card, start members at column head row
                member_y = col_head_y

            for member in col.get("members") or []:
                if member_y + ch > int(SLIDE_H) - int(Inches(0.12)):
                    break   # reached slide bottom — members beyond are in next slide
                from pptx.dml.color import RGBColor as _RC
                r, g, b = acc
                lighter = (min(255,r+40), min(255,g+40), min(255,b+40))
                _draw_card(slide, cx, member_y, member, lighter, card_w=cw)
                member_y += ch + int(CARD_V_GAP)

    # Chunk columns across slides
    chunks = [columns[i:i+MAX_COLS] for i in range(0, max(len(columns),1), MAX_COLS)]
    if not chunks:
        chunks = [[]]
    total = len(chunks)
    for idx, chunk in enumerate(chunks):
        sfx = f"({idx+1}/{total})" if total > 1 else ""
        _one_slide(chunk, sfx)


# ── Public API ────────────────────────────────────────────────────────────────

def build_pptx(company: str, industry: str,
               depts: list[dict[str, Any]]) -> bytes:
    """
    depts entries:
      { label, color, headcount, head: person|None,
        columns: [{head, members, accent_rgb}, …]
        head=None on a column means "continuation" — no column-head card drawn. }
    """
    prs = Presentation()
    prs.slide_width  = int(SLIDE_W)
    prs.slide_height = int(SLIDE_H)

    total = sum(d.get("headcount", 0) for d in depts)
    _make_cover(prs, company, total, len(depts), industry)

    for dept in depts:
        hc = dept.get("headcount", 0)
        if hc == 0 and not dept.get("head"):
            continue
        _make_dept_slides(
            prs,
            company        = company,
            dept_label     = dept.get("label", "Department"),
            dept_color_hex = dept.get("color", "#1e6eb8"),
            total_hc       = hc,
            dept_head      = dept.get("head"),
            columns        = dept.get("columns") or [],
        )

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
